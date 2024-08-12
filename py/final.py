import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
import os
import requests
import google.generativeai as genai
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound
import lancedb
import pyarrow as pa
import urllib.request as ur
from bs4 import BeautifulSoup
from gtts import gTTS
from io import BytesIO
import speech_recognition as sr
from googletrans import Translator
import time
import numpy as np
import pandas as pd
from langchain.schema import Document
from pytube import YouTube
import moviepy.editor as mp
import speech_recognition as sr
from typing import Optional, List
from PIL import Image, ImageEnhance

# Define a vector type for embeddings
def vector(dimension: int, value_type: pa.DataType = pa.float32()) -> pa.DataType:
    return pa.list_(value_type, dimension)

# Initialize the Translation API client
def gemini_api_embeddings(texts):
    return np.random.rand(len(texts), 768).tolist()

# LanceDB setup
table_name = "model5"
url = "D:\\lance"
try:
    db = lancedb.connect(url)
    schema = pa.schema([
        pa.field('id', pa.int32()),
        pa.field('file_name', pa.string()),
        pa.field('text', pa.string()),
        pa.field('embedding', vector(768))  # Use the vector function

    ])
    if table_name not in db.table_names():
        db.create_table(table_name, schema=schema)
except Exception as e:
    print(f"An exception occurred: {e}")

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

class GeminiAPIEmbeddings:
    def generate_embeddings(self, texts):
        return gemini_api_embeddings(texts)

    def ndims(self):
        return 768

gemini_func = GeminiAPIEmbeddings()

class Chunk:
    def __init__(self, id,file_name, text, embedding):
        self.id = id
        self.file_name = file_name
        self.text = text
        self.embedding = embedding

    def dict(self):
        return {
            "id": self.id,
            "file_name": self.file_name,
            "text": self.text,
            "embedding": self.embedding
        }

chunk_table = db.open_table(table_name)

# Define the rerank_hybrid function
def rerank_hybrid(query: str, vector_results: pa.Table, filters: list) -> pa.Table:
    # Convert to pandas DataFrame
    df = vector_results.to_pandas()
    
    # Calculate the number of words in each text entry
    df['num_words'] = df['text'].apply(lambda x: len(x.split()))
    
    # Apply filters using the calculated 'num_words'
    df = df.query("num_words > 150")
    
    # Convert back to PyArrow Table
    return pa.Table.from_pandas(df)


class ModifiedLinearReranker:
    def __init__(self, filters):
        self.filters = filters

    def rerank(self, query, vector_results):
        return rerank_hybrid(query, vector_results, self.filters)

# Function to extract text from PDF files
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf.seek(0)  # Ensure you're reading from the beginning of the file
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Function to extract text from PPTX files
def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        ppt.seek(0)  # Ensure you're reading from the beginning of the file
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text


    return text

# Function to extract text from a YouTube video URL
def process_link(url, lang='en'):
    """Extracts text from a YouTube video URL."""
    video_id = None
    if "youtube.com/watch?v=" in url:
        video_id = url.split("youtube.com/watch?v=")[1].split("&")[0]
    elif "youtu.be/" in url:
        video_id = url.split("youtu.be/")[1].split("?")[0]
    if video_id:
        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=[lang])
            text = " ".join([entry['text'] for entry in transcript])
            return text
        except TranscriptsDisabled:
            return "Transcripts are disabled for this video."
        except NoTranscriptFound:
            return "No transcript found for this video."
        except Exception as e:
            return f"Error retrieving transcript: {str(e)}"
    return "Invalid YouTube URL"

# Function to split text into chunks
def get_chunks_lance(raw_text,file_name):
    if not raw_text.strip():
        return
    chunk_size = 5000
    chunks = [raw_text[i:i+chunk_size] for i in range(0, len(raw_text), chunk_size)]
    for idx, chun in enumerate(chunks):
        chunk_embedding = gemini_func.generate_embeddings([chun])[0]
        chunk_instance = Chunk(id=idx,file_name=file_name, text=chun, embedding=chunk_embedding)
        chunk_table.add([chunk_instance.dict()])

def record_audio(lang='en-US'):
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        st.write("Listening...")
        audio = recognizer.listen(source)
        try:
            text = recognizer.recognize_google(audio, language=lang)
            return text
        except sr.UnknownValueError:
            return "Sorry, I did not understand that."
        except sr.RequestError as e:
            return f"Sorry, there was an error: {e}"

def text_to_audio(text, language):
    tts = gTTS(text, lang=language)
    audio_file = BytesIO()
    tts.write_to_fp(audio_file)
    audio_file.seek(0)
    return audio_file

def get_website_text(urlinput):
    try:
        response = ur.urlopen(urlinput)
        soup = BeautifulSoup(response.read(), 'html.parser')
        # Get all text content from the webpage
        text_content = soup.get_text(separator=' ')
        # Remove extra spaces
        cleaned_text_content = ' '.join(text_content.split())
        return cleaned_text_content
    except Exception as e:
        return f"Error retrieving website content: {str(e)}"

# Function to create a conversational AI chain
def get_conv():
    prompt_template = """
    You are a seminar taker teaching from a given document. Begin each explanation with an engaging and dynamic introduction such as and not exactly 'Hello newbie! Today, we are going to learn about...'. Make sure the introduction is generated on the spot and unique each time. Use simple language and make the lesson interesting and easy to understand. Provide detailed and thorough explanations.
    Context: {context}
    Question: {question}
    Answer (provide an answer which is relevant to the question mostly from the contents in the database,if no revelent answer is found for the question asked by the user generate the answer and also provide some examples for the asked question with some images relavant to the question):
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user input and generate responses

def fetch_image(query):
    search_url = "https://www.google.com/search?hl=en&tbm=isch&q=" + '+'.join(query.split())
    headers = {"User-Agent": "Mozilla/5.0"}
    response = requests.get(search_url, headers=headers)
    soup = BeautifulSoup(response.text, 'html.parser')
    img_tags = soup.find_all("img")
    
    for img_tag in img_tags[1:]:  # Start from the second image
        img_url = img_tag['src']
        response = requests.get(img_url)
        img = Image.open(BytesIO(response.content))
        
        # Skip images that are too small (likely to contain text or be low quality)
        if img.size[0] < 100 or img.size[1] < 100:
            continue
        
        # Skip images that are not in RGB mode (e.g., icons or other graphics)
        if img.mode != 'RGB':
            continue
        
        # Resize the image to a more reasonable size, using a high-quality filter
        resized_img = img.resize((300, 300), Image.LANCZOS)
        
        # Enhance the image clarity using ImageEnhance
        enhancer = ImageEnhance.Sharpness(resized_img)
        enhanced_img = enhancer.enhance(3)  # Increase sharpness
        
        return enhanced_img
    
    return None



def user_input(user_question,selected_doc):
    embeddings_model = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    tbl = db.open_table(table_name)
    tbl.create_fts_index(["text"], replace=True)
    question_embedding = embeddings_model.embed_query(user_question)

    #data=tbl.search((question_embedding, user_question.lower()), vector_column_name="embedding", query_type="hybrid")
    #print(data)
    data = tbl.search((question_embedding, user_question.lower()), vector_column_name="embedding",
                      query_type="hybrid").where(f"file_name = '{selected_doc}'")
    print(data)

    modified_reranker = ModifiedLinearReranker(filters=["dual-band"])
    reranked_results = modified_reranker.rerank(user_question, data)
    # Convert reranked results to documents
    reranked_chunks = reranked_results.to_pandas()['text'].tolist()
    documents = [Document(page_content=chunk) for chunk in reranked_chunks]
    # Use the reranked results in your conversational AI chain
    chain = get_conv()
    response = chain({"input_documents": documents, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

# Function to translate text to a desired language
def translate_text(text, target_language):
    translator = Translator()
    try:
        translated = translator.translate(text, dest=target_language)
        return translated.text
    except AttributeError:
        # Handle the specific case where the regex does not find a match
        time.sleep(2)  # Add a delay and try again
        translated = translator.translate(text, dest=target_language)
        return translated.text
    except Exception as e:
        # Handle other exceptions
        return f"Translation error: {str(e)}"

# Main function to run the Streamlit application
def main():
    images = []
    st.set_page_config(page_title="Chat with multiple PDFs, PPTs, and YouTube Videos")
    
    # Initialize session state variables
    if "user_question" not in st.session_state:
        st.session_state.user_question = ""
    if "state" not in st.session_state:
        st.session_state.state = "initial"
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "selected_language" not in st.session_state:
        st.session_state.selected_language = "en"  # Default language
    
    # Use custom CSS to set the background gradient and fix the text area
    st.markdown(
        f"""
        <style>
        .stApp {{
                background-color: #4158D0;
        background-image: linear-gradient(43deg, #4158D0 0%, #C850C0 46%, #FFCC70 100%);
        font-family: 'Arial', sans-serif;

        }}
        .fixed-bottom {{
            position: fixed;
            bottom: 0;
            left: 0;
            width: 100%;
            background-color: black;
            padding: 10px;
            z-index: 100;
        }}
        .gradient-text {{
            background: linear-gradient(to right, #ffffff, #f1f1f1);
        -webkit-background-clip: text;
            color: transparent;
            font-size: 48px;
            font-weight: bold;
        }}
        .sub-text {{
            color: blue;
            font-size: 36px;
        }}
        .matte-black {{
            background-color: #282828;
            color: white;
            padding: 10px;
            border-radius: 5px;
        }}
        </style>
        """,
        unsafe_allow_html=True
    )
    
    # Display the gradient header text
    st.markdown('<div class="gradient-text">Chat Bot</div>', unsafe_allow_html=True)
     
    with st.sidebar:
        st.title("Menu:")
        st.title("Scroll Down To Submit")
        if 'uploaded_files' not in st.session_state:
            st.session_state.uploaded_files = []

        file_list = [row['file_name'] for _, row in
                     chunk_table.to_pandas().drop_duplicates(subset='file_name').iterrows()]

        # Add a dropdown for selecting a file
        selected_file = st.selectbox("Select a file from the database", options=file_list)
        pdf_docs = st.file_uploader("Upload your PDF files", accept_multiple_files=True, type="pdf")
        ppt_docs = st.file_uploader("Upload your PPT files", accept_multiple_files=True, type="pptx")
        link = st.text_input("Enter a YouTube video link")
        if link:
            file_name=st.text_input("Enter a title for your link")
        language = st.selectbox("Select the transcript language", options=['en', 'es', 'fr', 'de', 'it', 'hi','ta'])
        urlinput = st.text_input("Enter a website URL")
        if urlinput:
            file_name=st.text_input("Enter a title for your link")
        if st.button("Submit"):
            with st.spinner("Processing.."):
                raw_text = ""
                if pdf_docs:
                    file_name = pdf_docs[0].name
                    st.write(f"Processing PDF: {file_name}")

                    raw_text += get_pdf_text(pdf_docs)

                if ppt_docs:
                    file_name = ppt_docs[0].name
                    st.write(f"Processing PPT: {file_name}")
                    raw_text += get_ppt_text(ppt_docs)
                if link:
                    video_text = process_link(link, lang=language)
                    st.write(f"Processing Link: {file_name}")
                    if "Error" not in video_text:
                        raw_text += video_text
                    else:
                        st.error(video_text)
                if urlinput:
                    website_text = get_website_text(urlinput)
                    raw_text += website_text
                get_chunks_lance(raw_text,file_name)
                st.session_state.state = "final"
                st.success("Done")

    # Display chat history
    if st.session_state.chat_history:
       for chat in st.session_state.chat_history:
           st.write(chat)

    # Text input for user question
    user_question = st.text_input("Ask a question", value=st.session_state.user_question, key="text_input")

    # Create a row for the "Record" and "New Input" buttons
    col1, col2 = st.columns([2, 0.5])
    with col1:
        if st.button("Record 🎙"):
            with st.spinner("Listening..."):
                # Capture audio and set the user question
                st.session_state.user_question = record_audio(lang=language)
                # Trigger a rerun to update the text input field
                st.rerun()
    with col2:
        pass

    # Add a text area at the bottom for user input like a chatbot
    if st.button("Send"):
        with st.spinner("Generating response..."):
            response = user_input(user_question,selected_file)
            st.session_state.chat_history.append(f"User: {user_question}")
            st.session_state.chat_history.append(f"Bot: {response}")
            st.write(f"User: {user_question}")
            st.write(f"Bot: {response}")
            if images:
                for img in images:
                    st.image(img, caption="Image from PDF")

            # Fetch and display relevant images from the internet, resized and enhanced
            relevant_img = fetch_image(user_question)
            if relevant_img:
                st.image(relevant_img, caption="Relevant Image from the Internet")
            #audio_file = text_to_audio(response, language)
            #st.audio(audio_file, format='audio/mp3')
            try:
                audio_file = text_to_audio(response, language)
                st.audio(audio_file, format='audio/mp3')
            except Exception as e:
                return f"Translation error"

    # Language selection for translation
    translation_language = st.selectbox("Translate to", options=['en', 'es', 'fr', 'de', 'it', 'hi','ta'], key="translate_lang")
    if st.button("Translate"):
        with st.spinner("Translating..."):
            response = user_input(user_question,selected_file)
            translated_text = translate_text(response, translation_language)
            st.write(f"Translated: {translated_text}")
            audio_file = text_to_audio(translated_text, translation_language)
            st.audio(audio_file, format='audio/mp3')
            

    # Determine which GIF to display based on the state
    if st.session_state.state == "initial":
        gif_path = "static/pic1.gif"
    elif st.session_state.state == "final":
        gif_path = "static/pic3.gif"

    # Display the appropriate GIF
    st.image(gif_path)

if __name__ == "__main__":
    main()